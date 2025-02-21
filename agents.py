from pptx import Presentation
import pptx
from swarm import Swarm, Agent
from openai import OpenAI
from run_demo_loop import run_demo_loop
import json

prs = Presentation()

def _new_slide(title: str, description: str):
    """
    Step 1: Create a new slide in the PowerPoint presentation.
    Step 2: Add a title to the slide.
    Step 3: Add a description to the slide.
    Step 4: Save the presentation file.
    
    Parameters:
    title (str): The main heading of the slide.
    description (str): The detailed content explaining the title.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    shapes.title.text = title
    shapes[1].text = description
    prs.save("./pptx-presentations/new.pptx")



# Modify the prs object to prevent overwriting the file every time.

router_agent = Agent(
    name="Router", 
    model="mistral-small:24b", 
    instructions="""

    Don't tell the user the steps just you follow them.


    Step 1: Receive the user's input.
    Step 2: Determine the nature of the request:
        - If the request is a general question, respond as a chatbot.
        - If the request is related to summarizing a book or making a presentation, forward the request to the Summarizer agent.
    Step 3: Ensure smooth communication with the Summarizer agent when necessary.
    Step 4: Provide accurate and meaningful responses to general questions.
    """
)

summarizer = Agent(
    name="Summarizer",
    model="mistral-small:24b",
    options={"temperature": 0},
    instructions="""
    
    Don't tell the user the steps just you follow them.

    Step 1: Greet the user and ask for the name of the book they want summarized.
    Step 2: Ask how many key points they want in the summary (default: 4).
    Step 3: Generate a summary based on the book and the specified number of key points.
    Step 4: Display the summary to the user in a readable format.
    Step 5: Ask the user if they want to add an overview, conclusion, or any other sections before proceeding.
    Step 6: Allow the user to modify or refine any parts of the summary if desired.
    Step 7: Once the user confirms the summary, format it for the presentation:
        - Titles: Generate key ideas separated by '^'.
        - Descriptions: Generate corresponding explanations separated by '^'.
        - Don't forgot to use "^" as a separator.
        - for a new line just a /n
    Step 8: Ask the user for final confirmation before calling make_presentation function.
    Step 9: If the user confirms, call make_presentation with the formatted titles and descriptions.
    Step 10: Inform the user once the presentation is successfully created.
    Step 11: Transfer control back to the Router agent for further interactions.
    """
)


def transfer_to_summarizer():
    """
    Step 1: This function redirects the request to the Summarizer agent.
    Step 2: Don't pass any arguments.
    """
    return summarizer

def transfer_to_router():
    """
    Step 1: This function redirects control back to the Router agent.
    Step 2: It ensures further user interaction if needed.
    """
    return router_agent

def make_presentation(titles, descriptions):
    """
    Step 1: Receive formatted titles and descriptions from the Summarizer agent.
    Step 2: Split the provided strings using the '^' separator.
    Step 3: Pair each title with its corresponding description.
    Step 4: Iterate over the title-description pairs and create slides accordingly.
    Step 5: Save the PowerPoint presentation file.
    Step 6: Return 'SUCCESS!' upon completion.
    
    Parameters:
    titles (str): A string containing slide titles separated by '^'.
    descriptions (str): A string containing slide descriptions separated by '^'.
    
    Example:
    titles = "Introduction^Main Idea^Conclusion"
    descriptions = "Brief overview^Detailed discussion^Final thoughts"
    """
    print("Making presentation...")

    titles = titles.split("^")
    descriptions = descriptions.split("^")
    presentation_values = dict(zip(titles, descriptions))
    for title, description in presentation_values.items():
        title = title.replace("#", "").replace("\n", "").replace("/", "")
        description = description.replace("#", "").replace("\n", "").replace("/", "")
        print("Creating slides...")
        print("Title:", title)
        print("Description:", description)
        _new_slide(title, description)
    return "SUCCESS!"

router_agent.functions.append(transfer_to_summarizer)
summarizer.functions = [transfer_to_router, make_presentation]

def _run_demo_loop(messages_starting, last_agent):
    if last_agent is None: 
        last_agent = router_agent
    ollama_client = OpenAI(
        base_url="http://localhost:11434/v1",        
        api_key="ollama"            
    )
    client = Swarm(client=ollama_client)

    messages = client.run(
        agent=last_agent,
        messages=messages_starting,
    )
    content = []
    for message in messages.messages:
        print("print message...")
        print(message)
        if message["role"] != "assistant":
            continue

        # Print response, if any
        if message["content"]:
            content = message["content"]

        # Print tool calls in purple, if any
        tool_calls = message.get("tool_calls") or []
        if len(tool_calls) > 1:
            print("print tool call...")
        for tool_call in tool_calls:
            f = tool_call["function"]
            name, args = f["name"], f["arguments"]
            arg_str = json.dumps(json.loads(args)).replace(":", "=")
            tool_call = (f"{name}({arg_str[1:-1]})")
    last_agent = messages.agent
    role =  message["role"]
    return content, role, last_agent

run_demo_loop(router_agent)
